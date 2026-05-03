import logging
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation

from parsers.base import BaseParser, ParsedChunk

logger = logging.getLogger(__name__)


class PptxParser(BaseParser):
    def parse(self, file_path: Path, base_dir: Path) -> list[ParsedChunk]:
        try:
            prs = Presentation(str(file_path))
        except Exception:
            logger.warning("Failed to parse %s, skipping.", file_path)
            return []

        last_modified = datetime.fromtimestamp(
            file_path.stat().st_mtime, tz=timezone.utc
        ).isoformat()

        chunks: list[ParsedChunk] = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        texts.append(text)

            if texts:
                chunks.append(
                    self._make_chunk(
                        "\n".join(texts),
                        file_path,
                        base_dir,
                        "pptx",
                        slide_num,
                        last_modified,
                    )
                )

        return chunks
